from urllib.parse import quote
from io import BytesIO

from django.conf import settings
from django.core.files.base import ContentFile
from django.db import connection
from django.db.models import Count, Q
from rest_framework import permissions, status, viewsets
from rest_framework.decorators import action
from rest_framework.exceptions import PermissionDenied, ValidationError
from rest_framework.response import Response

from .access import (
    can_edit_document,
    can_edit_folder,
    can_share_document,
    can_share_folder,
    can_view_document,
    can_view_folder,
    document_queryset_for_user,
    expand_folder_descendants,
    folder_queryset_for_user,
    user_can_manage_ged_share,
)
from .models import Document, DocumentFolder, GedShare
from .notifications import (
    notify_ged_share_created,
    notify_ged_share_revoked,
    notify_ged_share_role_changed,
    _item_label,
)
from .onlyoffice import (
    build_onlyoffice_config,
    extension_from_document,
    jwt_sign_onlyoffice_config,
    office_file_token,
    onlyoffice_document_type_for_ext,
)
from .serializers import DocumentFolderSerializer, DocumentSerializer, GedShareSerializer


def _build_empty_office_file(kind: str) -> tuple[bytes, str, str]:
    kind = (kind or "").strip().lower()
    if kind == "docx":
        try:
            from docx import Document as DocxDocument
        except ImportError as exc:
            raise ValidationError({"detail": "Support DOCX indisponible (python-docx non installe)."}) from exc
        doc = DocxDocument()
        doc.add_paragraph("")
        buf = BytesIO()
        doc.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "docx",
        )
    if kind == "xlsx":
        try:
            from openpyxl import Workbook
        except ImportError as exc:
            raise ValidationError({"detail": "Support XLSX indisponible (openpyxl non installe)."}) from exc
        wb = Workbook()
        ws = wb.active
        ws.title = "Feuille1"
        buf = BytesIO()
        wb.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "xlsx",
        )
    if kind == "pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ValidationError({"detail": "Support PPTX indisponible (python-pptx non installe)."}) from exc
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        buf = BytesIO()
        prs.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx",
        )
    raise ValidationError({"kind": "Type invalide. Utilisez docx, xlsx ou pptx."})


class DocumentFolderViewSet(viewsets.ModelViewSet):
    permission_classes = [permissions.IsAuthenticated]
    serializer_class = DocumentFolderSerializer

    def get_queryset(self):
        user = self.request.user
        qs = folder_queryset_for_user(user).annotate(
            _ged_child_count=Count("children", distinct=True),
            _ged_doc_count=Count("documents", distinct=True),
        ).select_related("parent", "created_by").order_by("name")
        parent = self.request.query_params.get("parent")
        if parent == "":
            qs = qs.filter(parent__isnull=True)
        elif parent:
            qs = qs.filter(parent_id=parent)
        search = (self.request.query_params.get("search") or "").strip()
        if search:
            qs = qs.filter(name__icontains=search)
        if self.request.query_params.get("scope") == "shared":
            direct = set(
                GedShare.objects.filter(shared_with=user, folder_id__isnull=False).values_list(
                    "folder_id", flat=True
                )
            )
            expanded = expand_folder_descendants(direct) | direct
            qs = qs.filter(id__in=expanded)
        return qs

    def get_serializer_context(self):
        ctx = super().get_serializer_context()
        ctx["request"] = self.request
        return ctx

    def perform_create(self, serializer):
        user = self.request.user
        parent = serializer.validated_data.get("parent")
        if parent and not can_edit_folder(user, parent):
            raise PermissionDenied("Creation de dossier non autorisee ici.")
        serializer.save(created_by=user)

    def perform_update(self, serializer):
        user = self.request.user
        instance = self.get_object()
        if not can_edit_folder(user, instance):
            raise PermissionDenied("Modification non autorisee.")
        serializer.save()

    def perform_destroy(self, instance):
        user = self.request.user
        if not can_edit_folder(user, instance):
            raise PermissionDenied("Suppression non autorisee.")
        has_content = instance.children.exists() or instance.documents.exists()
        if has_content and not (user.is_staff or user.is_superuser):
            raise PermissionDenied(
                "Seuls les administrateurs peuvent supprimer un dossier contenant des sous-dossiers ou des documents."
            )
        instance.delete()


class DocumentViewSet(viewsets.ModelViewSet):
    permission_classes = [permissions.IsAuthenticated]
    serializer_class = DocumentSerializer

    def get_queryset(self):
        user = self.request.user
        qs = document_queryset_for_user(user)
        folder_id = self.request.query_params.get("folder")
        if folder_id == "":
            qs = qs.filter(folder__isnull=True)
        elif folder_id:
            qs = qs.filter(folder_id=folder_id)
        search = (self.request.query_params.get("search") or "").strip()
        if search:
            qs = qs.filter(
                Q(title__icontains=search)
                | Q(original_filename__icontains=search)
                | Q(description__icontains=search)
            )
        if self.request.query_params.get("scope") == "shared":
            direct_f = set(
                GedShare.objects.filter(shared_with=user, folder_id__isnull=False).values_list(
                    "folder_id", flat=True
                )
            )
            folder_ids = expand_folder_descendants(direct_f) | direct_f
            qs = qs.filter(
                Q(id__in=GedShare.objects.filter(shared_with=user, document_id__isnull=False).values_list(
                    "document_id", flat=True
                ))
                | Q(folder_id__in=folder_ids)
            )
        return qs

    def get_serializer_context(self):
        ctx = super().get_serializer_context()
        ctx["request"] = self.request
        return ctx

    def perform_create(self, serializer):
        user = self.request.user
        folder = serializer.validated_data.get("folder")
        if folder and not can_edit_folder(user, folder):
            raise PermissionDenied("Depot non autorise dans ce dossier.")
        serializer.save(uploaded_by=user)

    def perform_update(self, serializer):
        user = self.request.user
        instance = self.get_object()
        if not can_edit_document(user, instance):
            raise PermissionDenied("Modification non autorisee.")
        serializer.save()

    def perform_destroy(self, instance):
        user = self.request.user
        if not can_edit_document(user, instance):
            raise PermissionDenied("Suppression non autorisee.")
        instance.delete()

    @action(detail=False, methods=["post"], url_path="create-office")
    def create_office(self, request):
        kind = (request.data.get("kind") or "").strip().lower()
        title = (request.data.get("title") or "").strip()
        folder_id = request.data.get("folder")
        folder = None
        if folder_id:
            folder = DocumentFolder.objects.filter(pk=folder_id).first()
            if not folder or not can_edit_folder(request.user, folder):
                raise PermissionDenied("Creation non autorisee dans ce dossier.")

        content, mime, ext = _build_empty_office_file(kind)
        base_title = title or f"Nouveau {ext.upper()}"
        filename = f"{base_title}.{ext}".replace("/", "-").replace("\\", "-").strip()

        doc = Document(
            title=base_title,
            description="",
            folder=folder,
            original_filename=filename,
            mime_type=mime,
            file_size=len(content),
            uploaded_by=request.user,
        )
        doc.file.save(filename, ContentFile(content), save=False)
        doc.save()

        serializer = self.get_serializer(doc)
        return Response(serializer.data, status=status.HTTP_201_CREATED)

    @action(detail=True, methods=["get"], url_path="preview-url")
    def preview_url(self, request, pk=None):
        """URL signee pour iframe / video / PDF (meme jeton que ONLYOFFICE)."""
        doc = self.get_object()
        if not can_view_document(request.user, doc):
            raise PermissionDenied()
        schema = connection.schema_name
        t = office_file_token(str(doc.pk), schema)
        path = f"/api/ged/office-file/?t={quote(t, safe='')}"
        return Response({"url": request.build_absolute_uri(path)})

    @action(detail=True, methods=["get"], url_path="onlyoffice-config")
    def onlyoffice_config(self, request, pk=None):
        doc = self.get_object()
        if not can_view_document(request.user, doc):
            raise PermissionDenied()
        if not (settings.ONLYOFFICE_JWT_SECRET or "").strip():
            return Response(
                {"detail": "ONLYOFFICE non configure (ONLYOFFICE_JWT_SECRET)."},
                status=status.HTTP_503_SERVICE_UNAVAILABLE,
            )
        ext = extension_from_document(doc)
        if onlyoffice_document_type_for_ext(ext) is None:
            return Response(
                {"detail": "Format non pris en charge par ONLYOFFICE pour ce fichier."},
                status=status.HTTP_400_BAD_REQUEST,
            )
        internal = (settings.OFFICE_FILE_DOWNLOAD_BASE or "").strip()
        if not internal:
            return Response(
                {"detail": "OFFICE_FILE_DOWNLOAD_BASE non configure."},
                status=status.HTTP_503_SERVICE_UNAVAILABLE,
            )
        schema = connection.schema_name
        editable = can_edit_document(request.user, doc)
        try:
            config = build_onlyoffice_config(doc, schema, request.user, internal, can_edit=editable)
            token = jwt_sign_onlyoffice_config(config)
        except ValueError as exc:
            if str(exc) == "missing_jwt_secret":
                return Response(
                    {"detail": "ONLYOFFICE non configure."},
                    status=status.HTTP_503_SERVICE_UNAVAILABLE,
                )
            raise
        return Response(
            {
                "documentServerUrl": settings.ONLYOFFICE_DOCUMENT_SERVER_URL.rstrip("/"),
                "config": config,
                "token": token,
            }
        )


class GedShareViewSet(viewsets.ModelViewSet):
    permission_classes = [permissions.IsAuthenticated]
    serializer_class = GedShareSerializer
    http_method_names = ["get", "post", "patch", "delete", "head", "options"]

    def get_queryset(self):
        base = GedShare.objects.select_related("folder", "document", "shared_with", "shared_by")
        user = self.request.user

        # Liste : filtre par dossier ou document (URL ?folder= / ?document=)
        if self.action == "list":
            fid = self.request.query_params.get("folder")
            did = self.request.query_params.get("document")
            if fid:
                folder = DocumentFolder.objects.filter(pk=fid).first()
                if not folder or not can_view_folder(user, folder):
                    return GedShare.objects.none()
                ancestor_ids: list = []
                walk = folder.parent
                while walk:
                    ancestor_ids.append(walk.pk)
                    walk = walk.parent
                q = Q(folder_id=fid)
                if ancestor_ids:
                    q |= Q(folder_id__in=ancestor_ids)
                return base.filter(q).order_by("-created_at")
            if did:
                doc = Document.objects.filter(pk=did).select_related("folder").first()
                if not doc or not can_view_document(user, doc):
                    return GedShare.objects.none()
                folder_ids: list = []
                walk = doc.folder
                while walk:
                    folder_ids.append(walk.pk)
                    walk = walk.parent
                q = Q(document_id=did)
                if folder_ids:
                    q |= Q(folder_id__in=folder_ids)
                return base.filter(q).order_by("-created_at")
            return GedShare.objects.none()

        # PATCH / DELETE / GET detail : pas de ?folder / ?document dans l'URL — il faut pouvoir
        # resoudre l'objet par pk (sinon get_queryset() vide => 404 permanent).
        if user.is_staff or user.is_superuser:
            return base.order_by("-created_at")
        visible_folder_ids = folder_queryset_for_user(user).values_list("id", flat=True)
        visible_document_ids = document_queryset_for_user(user).values_list("id", flat=True)
        return (
            base.filter(
                Q(shared_by=user)
                | Q(shared_with=user)
                | Q(folder_id__in=visible_folder_ids)
                | Q(document_id__in=visible_document_ids)
            )
            .order_by("-created_at")
            .distinct()
        )

    def get_serializer_context(self):
        ctx = super().get_serializer_context()
        ctx["request"] = self.request
        ctx["ged_share_list_folder_id"] = self.request.query_params.get("folder")
        did = self.request.query_params.get("document")
        ctx["ged_share_list_document_id"] = did
        ctx["ged_share_list_document_folder_id"] = None
        if did:
            doc = Document.objects.filter(pk=did).only("folder_id").first()
            if doc and doc.folder_id:
                ctx["ged_share_list_document_folder_id"] = str(doc.folder_id)
        return ctx

    def perform_create(self, serializer):
        user = self.request.user
        folder = serializer.validated_data.get("folder")
        document = serializer.validated_data.get("document")
        if folder:
            if not can_view_folder(user, folder):
                raise PermissionDenied("Dossier introuvable.")
            if not can_share_folder(user, folder):
                raise PermissionDenied("Partage non autorise pour ce dossier.")
        elif document:
            if not can_view_document(user, document):
                raise PermissionDenied("Document introuvable.")
            if not can_share_document(user, document):
                raise PermissionDenied("Partage non autorise pour ce document.")
        else:
            raise ValidationError({"detail": "Dossier ou document requis."})
        share = serializer.save(shared_by=user)
        try:
            notify_ged_share_created(tenant=getattr(self.request, "tenant", None), share=share, actor=user)
        except Exception:
            pass

    def perform_update(self, serializer):
        user = self.request.user
        instance = serializer.instance
        if not user_can_manage_ged_share(user, instance):
            raise PermissionDenied("Modification non autorisee.")
        old_role = instance.role
        share = serializer.save()
        if share.role != old_role:
            try:
                notify_ged_share_role_changed(
                    tenant=getattr(self.request, "tenant", None),
                    share=share,
                    actor=user,
                    old_role=old_role,
                )
            except Exception:
                pass

    def perform_destroy(self, instance):
        user = self.request.user
        if not user_can_manage_ged_share(user, instance):
            raise PermissionDenied("Suppression non autorisee.")
        tenant = getattr(self.request, "tenant", None)
        recipient = instance.shared_with
        item_label = _item_label(instance)
        try:
            notify_ged_share_revoked(tenant=tenant, recipient=recipient, item_label=item_label, actor=user)
        except Exception:
            pass
        instance.delete()
